"""
ProductIQ — Report Generation Tools
PDF (WeasyPrint + Jinja2), PPTX (python-pptx), Supabase upload, RFQ generator.
"""

import json
import os
import tempfile
from pathlib import Path
from crewai.tools import BaseTool
from database import get_supabase
import structlog

logger = structlog.get_logger()

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"
REPORTS_TMP = Path(tempfile.gettempdir()) / "productiq_reports"
REPORTS_TMP.mkdir(exist_ok=True)


# ── PDF Generator ─────────────────────────────────────────────────────────────

class PDFGeneratorTool(BaseTool):
    name: str = "PDF Report Generator"
    description: str = (
        "Generates a branded PDF executive report from structured report data. "
        "Input: report_data_json (dict with all agent output data), "
        "output_filename (e.g. 'report_abc123.pdf'), watermarked (bool as string 'true'/'false')."
    )

    def _run(self, report_data_json: str, output_filename: str, watermarked: str = "false") -> str:
        try:
            from weasyprint import HTML, CSS
            from jinja2 import Environment, FileSystemLoader

            data = json.loads(report_data_json) if isinstance(report_data_json, str) else report_data_json
            is_watermarked = str(watermarked).lower() == "true"

            # Load Jinja2 template
            env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)))
            try:
                template = env.get_template("report.html")
            except Exception:
                # Fallback: generate minimal HTML if template is missing
                template_str = _get_fallback_html_template()
                from jinja2 import Template
                template = Template(template_str)

            html_content = template.render(**data, watermarked=is_watermarked)

            output_path = str(REPORTS_TMP / output_filename)
            HTML(
                string=html_content,
                base_url=str(TEMPLATES_DIR),
            ).write_pdf(output_path)

            size = os.path.getsize(output_path)
            logger.info("PDF generated", path=output_path, size=size)
            return json.dumps({
                "path": output_path,
                "size_bytes": size,
                "filename": output_filename,
                "watermarked": is_watermarked,
            })

        except ImportError:
            return json.dumps({"error": "WeasyPrint not installed. Run: pip install weasyprint"})
        except Exception as e:
            logger.error("PDF generation error", error=str(e))
            return json.dumps({"error": str(e)})


# ── PPTX Generator ────────────────────────────────────────────────────────────

class PPTXGeneratorTool(BaseTool):
    name: str = "PowerPoint Report Generator"
    description: str = (
        "Generates a branded PPTX presentation from structured report data using python-pptx. "
        "Input: report_data_json (dict), output_filename (e.g. 'report_abc123.pptx')."
    )

    def _run(self, report_data_json: str, output_filename: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            data = json.loads(report_data_json) if isinstance(report_data_json, str) else report_data_json

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Colour palette
            DARK = RGBColor(0x0D, 0x0D, 0x14)
            ACCENT = RGBColor(0x6C, 0x63, 0xFF)
            WHITE = RGBColor(0xFF, 0xFF, 0xFF)
            LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF8)
            TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)

            def add_slide(layout_idx=6):
                slide_layout = prs.slide_layouts[layout_idx]
                return prs.slides.add_slide(slide_layout)

            def set_bg(slide, color: RGBColor):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = color

            def add_text_box(slide, text: str, left, top, width, height,
                             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = align
                run = p.add_run()
                run.text = str(text)
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = color
                return txBox

            def add_rect(slide, left, top, width, height, color: RGBColor):
                shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
                shape.line.fill.background()
                return shape

            category = data.get("category", "Product Category")
            brand = data.get("brand_name", "Brand")

            # ── Slide 1: Cover ────────────────────────────────────────────────
            slide = add_slide()
            set_bg(slide, DARK)
            add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), ACCENT)
            add_text_box(slide, "ProductIQ", Inches(0.6), Inches(0.5), Inches(4), Inches(0.6),
                         font_size=14, color=ACCENT)
            add_text_box(slide, f"Product Intelligence Report", Inches(0.6), Inches(1.5), Inches(11), Inches(1.2),
                         font_size=40, bold=True, color=WHITE)
            add_text_box(slide, f"{category} | {brand}", Inches(0.6), Inches(2.9), Inches(11), Inches(0.8),
                         font_size=24, color=RGBColor(0xCC, 0xCC, 0xDD))
            add_text_box(slide, f"AI-Generated Market Intelligence  •  {data.get('date', 'April 2026')}",
                         Inches(0.6), Inches(6.5), Inches(11), Inches(0.5), font_size=12,
                         color=RGBColor(0x88, 0x88, 0x99))

            # ── Slide 2: Executive Summary ─────────────────────────────────────
            slide = add_slide()
            set_bg(slide, WHITE)
            add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT)
            add_text_box(slide, "Executive Summary", Inches(0.4), Inches(0.3), Inches(11), Inches(0.7),
                         font_size=28, bold=True, color=TEXT_DARK)

            top_insights = data.get("top_insights", data.get("insights", []))[:5]
            for i, ins in enumerate(top_insights):
                y = 1.2 + i * 1.1
                add_rect(slide, Inches(0.4), Inches(y), Inches(11.8), Inches(0.9), LIGHT_GRAY)
                title = ins.get("title", f"Insight {i+1}") if isinstance(ins, dict) else str(ins)
                body = ins.get("body", "")[:120] if isinstance(ins, dict) else ""
                add_text_box(slide, f"  {i+1}. {title}", Inches(0.4), Inches(y+0.05), Inches(11.8), Inches(0.4),
                             font_size=14, bold=True, color=TEXT_DARK)
                if body:
                    add_text_box(slide, f"  {body}...", Inches(0.4), Inches(y+0.45), Inches(11.8), Inches(0.35),
                                 font_size=11, color=RGBColor(0x44, 0x44, 0x55))

            # ── Slide 3: Market Overview ───────────────────────────────────────
            slide = add_slide()
            set_bg(slide, WHITE)
            add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT)
            add_text_box(slide, "Market Overview", Inches(0.4), Inches(0.3), Inches(11), Inches(0.7),
                         font_size=28, bold=True, color=TEXT_DARK)

            products = data.get("products", [])
            total_prods = len(products)
            prices = [p.get("price_inr") for p in products if p.get("price_inr")]
            avg_price = round(sum(prices) / len(prices), 0) if prices else 0
            ratings = [p.get("rating") for p in products if p.get("rating")]
            avg_rating = round(sum(ratings) / len(ratings), 1) if ratings else 0

            stats = [
                ("Products Analysed", f"{total_prods:,}"),
                ("Price Range", f"₹{min(prices, default=0):,.0f}–₹{max(prices, default=0):,.0f}"),
                ("Avg Market Price", f"₹{avg_price:,.0f}"),
                ("Avg Rating", f"{avg_rating} / 5.0"),
            ]
            for j, (label, value) in enumerate(stats):
                x = 0.4 + j * 3.1
                add_rect(slide, Inches(x), Inches(1.4), Inches(2.8), Inches(1.8), LIGHT_GRAY)
                add_text_box(slide, value, Inches(x + 0.15), Inches(1.6), Inches(2.5), Inches(0.8),
                             font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
                add_text_box(slide, label, Inches(x + 0.1), Inches(2.3), Inches(2.6), Inches(0.5),
                             font_size=11, color=RGBColor(0x66, 0x66, 0x77), align=PP_ALIGN.CENTER)

            # ── Slide 4: Consumer Intelligence ────────────────────────────────
            slide = add_slide()
            set_bg(slide, WHITE)
            add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT)
            add_text_box(slide, "Consumer Intelligence", Inches(0.4), Inches(0.3), Inches(11), Inches(0.7),
                         font_size=28, bold=True, color=TEXT_DARK)

            clusters = data.get("review_clusters", [])
            pain_points = [c for c in clusters if c.get("topic_type") == "pain_point"][:4]
            add_text_box(slide, "Top Pain Points:", Inches(0.4), Inches(1.2), Inches(6), Inches(0.5),
                         font_size=16, bold=True, color=TEXT_DARK)
            for i, pp in enumerate(pain_points):
                add_text_box(slide,
                             f"• {pp.get('topic_label', '')} ({pp.get('review_count', 0)} reviews)",
                             Inches(0.5), Inches(1.8 + i * 0.6), Inches(6), Inches(0.5),
                             font_size=13, color=TEXT_DARK)

            # ── Slide 5: Competitive Landscape ────────────────────────────────
            slide = add_slide()
            set_bg(slide, WHITE)
            add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT)
            add_text_box(slide, "Competitive Landscape", Inches(0.4), Inches(0.3), Inches(11), Inches(0.7),
                         font_size=28, bold=True, color=TEXT_DARK)

            competitors = data.get("competitors", [])[:6]
            if competitors:
                headers = ["Brand", "Price (₹)", "Rating", "Key Strength"]
                col_widths = [2.5, 1.5, 1.2, 5.3]
                x_positions = [0.4, 2.95, 4.5, 5.75]

                # Header row
                add_rect(slide, Inches(0.4), Inches(1.3), Inches(11.5), Inches(0.4), ACCENT)
                for hi, (h, w) in enumerate(zip(headers, col_widths)):
                    add_text_box(slide, h, Inches(x_positions[hi] + 0.05), Inches(1.33),
                                 Inches(w - 0.1), Inches(0.35), font_size=11, bold=True, color=WHITE)

                for i, comp in enumerate(competitors):
                    y = 1.75 + i * 0.6
                    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
                    add_rect(slide, Inches(0.4), Inches(y), Inches(11.5), Inches(0.5), bg)
                    row_vals = [
                        comp.get("brand_name", "")[:25],
                        f"₹{comp.get('price_inr', 0):,.0f}",
                        str(comp.get("rating", "")),
                        ", ".join(comp.get("key_strengths", ["—"])[:2])[:80],
                    ]
                    for hi, (val, w) in enumerate(zip(row_vals, col_widths)):
                        add_text_box(slide, str(val), Inches(x_positions[hi] + 0.05), Inches(y + 0.08),
                                     Inches(w - 0.1), Inches(0.35), font_size=11, color=TEXT_DARK)

            # ── Slide 6: Product Concepts ─────────────────────────────────────
            concepts = data.get("product_concepts", [])
            for ci, concept in enumerate(concepts[:3]):
                slide = add_slide()
                set_bg(slide, DARK)
                add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), ACCENT)
                add_text_box(slide, f"Product Concept {ci+1}", Inches(0.6), Inches(0.3),
                             Inches(11), Inches(0.5), font_size=14, color=ACCENT)
                add_text_box(slide, concept.get("concept_name", f"Concept {ci+1}"),
                             Inches(0.6), Inches(0.9), Inches(11), Inches(0.9),
                             font_size=34, bold=True, color=WHITE)
                add_text_box(slide, concept.get("tagline", ""),
                             Inches(0.6), Inches(1.8), Inches(11), Inches(0.5),
                             font_size=18, color=RGBColor(0xBB, 0xBB, 0xCC))
                add_text_box(slide, f"Target: {concept.get('target_persona', '')}",
                             Inches(0.6), Inches(2.5), Inches(6), Inches(0.4),
                             font_size=12, color=RGBColor(0x99, 0x99, 0xAA))
                add_text_box(slide, f"Price Point: ₹{concept.get('suggested_price_inr', '')}",
                             Inches(0.6), Inches(2.95), Inches(6), Inches(0.4),
                             font_size=12, color=RGBColor(0x99, 0x99, 0xAA))
                add_text_box(slide, f"Validation Score: {concept.get('validation_score', '')} / 100",
                             Inches(8), Inches(2.95), Inches(4), Inches(0.4),
                             font_size=20, bold=True, color=ACCENT)

                features = concept.get("key_features", [])[:5]
                add_text_box(slide, "Key Features:", Inches(0.6), Inches(3.6),
                             Inches(12), Inches(0.4), font_size=13, bold=True, color=WHITE)
                for fi, feat in enumerate(features):
                    add_text_box(slide, f"  ›  {feat}", Inches(0.6), Inches(4.1 + fi * 0.45),
                                 Inches(12), Inches(0.4), font_size=12, color=RGBColor(0xBB, 0xBB, 0xCC))

            # ── Slide: GTM Strategy ───────────────────────────────────────────
            gtm = data.get("gtm_plans", [{}])[0] if data.get("gtm_plans") else {}
            slide = add_slide()
            set_bg(slide, WHITE)
            add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT)
            add_text_box(slide, "Go-To-Market Strategy", Inches(0.4), Inches(0.3),
                         Inches(11), Inches(0.7), font_size=28, bold=True, color=TEXT_DARK)

            messaging = gtm.get("messaging_framework", {})
            if messaging:
                hero = messaging.get("hero_message", "") or messaging.get("hero_message", "")
                add_text_box(slide, f"Hero Message: \"{hero}\"",
                             Inches(0.4), Inches(1.2), Inches(11.5), Inches(0.6),
                             font_size=16, bold=True, color=ACCENT)

            channels = gtm.get("launch_channels", [])[:5]
            add_text_box(slide, "Launch Channels:", Inches(0.4), Inches(2.0),
                         Inches(6), Inches(0.4), font_size=14, bold=True, color=TEXT_DARK)
            for ci_idx, ch in enumerate(channels):
                ch_text = ch if isinstance(ch, str) else ch.get("channel", str(ch))
                add_text_box(slide, f"  {ci_idx+1}. {ch_text}",
                             Inches(0.5), Inches(2.5 + ci_idx * 0.55), Inches(6), Inches(0.45),
                             font_size=13, color=TEXT_DARK)

            # ── Final slide ───────────────────────────────────────────────────
            slide = add_slide()
            set_bg(slide, DARK)
            add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), ACCENT)
            add_text_box(slide, "Thank You", Inches(1), Inches(2),
                         Inches(11), Inches(1.5), font_size=52, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER)
            add_text_box(slide, "Generated by ProductIQ — productiq.in",
                         Inches(1), Inches(4), Inches(11), Inches(0.5),
                         font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)

            output_path = str(REPORTS_TMP / output_filename)
            prs.save(output_path)
            size = os.path.getsize(output_path)
            slide_count = len(prs.slides)

            logger.info("PPTX generated", path=output_path, slides=slide_count, size=size)
            return json.dumps({
                "path": output_path,
                "slide_count": slide_count,
                "size_bytes": size,
                "filename": output_filename,
            })

        except ImportError:
            return json.dumps({"error": "python-pptx not installed. Run: pip install python-pptx"})
        except Exception as e:
            logger.error("PPTX generation error", error=str(e))
            return json.dumps({"error": str(e)})


# ── Supabase Upload Tool ──────────────────────────────────────────────────────

class SupabaseUploadTool(BaseTool):
    name: str = "Supabase Storage Uploader"
    description: str = (
        "Uploads a file to Supabase Storage and returns a signed URL valid for 7 days. "
        "Input: file_path (absolute local path), bucket (storage bucket name, e.g. 'reports'), "
        "storage_path (path inside the bucket, e.g. 'user_id/run_id/report.pdf')."
    )

    def _run(self, file_path: str, bucket: str, storage_path: str) -> str:
        try:
            db = get_supabase()

            if not os.path.exists(file_path):
                return json.dumps({"error": f"File not found: {file_path}"})

            content_type = "application/pdf" if file_path.endswith(".pdf") else (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                if file_path.endswith(".pptx") else "application/octet-stream"
            )

            with open(file_path, "rb") as f:
                db.storage.from_(bucket).upload(
                    path=storage_path,
                    file=f,
                    file_options={"content-type": content_type, "upsert": "true"},
                )

            url_response = db.storage.from_(bucket).create_signed_url(storage_path, 604800)
            signed_url = url_response.get("signedURL") or url_response.get("signed_url", "")

            # Clean up temp file
            try:
                os.remove(file_path)
            except Exception:
                pass

            logger.info("File uploaded to Supabase Storage", storage_path=storage_path, url=signed_url[:60])
            return json.dumps({
                "signed_url": signed_url,
                "storage_path": storage_path,
                "bucket": bucket,
                "signed_url_expiry_seconds": 604800,
            })

        except Exception as e:
            logger.error("Supabase upload error", file_path=file_path, error=str(e))
            return json.dumps({"error": str(e)})


# ── Report Record Store Tool ──────────────────────────────────────────────────

class ReportStoreTool(BaseTool):
    name: str = "Supabase Report Record Storage"
    description: str = (
        "Stores the final report record with URLs into the 'reports' table. "
        "Input: run_id (string UUID), user_id (string UUID), pdf_url (string), pptx_url (string, optional)."
    )

    def _run(self, run_id: str, user_id: str, pdf_url: str, pptx_url: str = "") -> str:
        try:
            db = get_supabase()
            record = {
                "run_id": run_id,
                "user_id": user_id,
                "pdf_url": pdf_url,
                "pptx_url": pptx_url or None,
            }
            db.table("reports").insert([record]).execute()
            logger.info("Report record stored", run_id=run_id)
            return json.dumps({"success": True, "stored_record": record})
        except Exception as e:
            logger.error("Report store error", error=str(e))
            return json.dumps({"error": str(e)})


# ── RFQ Generator Tool ────────────────────────────────────────────────────────

class RFQGeneratorTool(BaseTool):
    name: str = "RFQ PDF Generator"
    description: str = (
        "Generates a Request for Quotation PDF for a manufacturer/supplier. "
        "Input: supplier_json (supplier dict), product_concept_json (concept dict)."
    )

    def _run(self, supplier_json: str, product_concept_json: str) -> str:
        try:
            from jinja2 import Environment, FileSystemLoader, Template

            supplier = json.loads(supplier_json) if isinstance(supplier_json, str) else supplier_json
            concept = json.loads(product_concept_json) if isinstance(product_concept_json, str) else product_concept_json

            try:
                env = Environment(loader=FileSystemLoader(str(TEMPLATES_DIR)))
                template = env.get_template("rfq.html")
            except Exception:
                template = Template(_get_rfq_fallback_template())

            html = template.render(supplier=supplier, concept=concept)

            company_name = str(supplier.get("company_name", "supplier")).replace(" ", "_")[:30]
            filename = f"rfq_{company_name}.pdf"
            output_path = str(REPORTS_TMP / filename)

            try:
                from weasyprint import HTML
                HTML(string=html).write_pdf(output_path)
                size = os.path.getsize(output_path)
                return json.dumps({
                    "rfq_path": output_path,
                    "filename": filename,
                    "size_bytes": size,
                    "supplier": supplier.get("company_name"),
                })
            except ImportError:
                # Save as HTML if WeasyPrint not available
                html_path = output_path.replace(".pdf", ".html")
                with open(html_path, "w", encoding="utf-8") as f:
                    f.write(html)
                return json.dumps({"rfq_path": html_path, "warning": "WeasyPrint unavailable — saved as HTML"})

        except Exception as e:
            logger.error("RFQ generation error", error=str(e))
            return json.dumps({"error": str(e)})


# ── Fallback HTML templates ───────────────────────────────────────────────────

def _get_fallback_html_template() -> str:
    return """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<title>ProductIQ Report — {{ category }}</title>
<style>
  @page { size: A4; margin: 2cm; }
  body { font-family: 'Helvetica Neue', Arial, sans-serif; color: #1a1a2e; line-height: 1.6; }
  h1 { color: #6c63ff; font-size: 2.5em; margin-bottom: 0.2em; }
  h2 { color: #1a1a2e; border-bottom: 2px solid #6c63ff; padding-bottom: 0.3em; margin-top: 2em; }
  .cover { text-align: center; padding: 4cm 2cm; background: #0d0d14; color: white; min-height: 20cm; }
  .cover h1 { color: #6c63ff; font-size: 3em; }
  .cover .sub { color: #ccccdd; font-size: 1.4em; margin-top: 1em; }
  .stat-grid { display: grid; grid-template-columns: repeat(4, 1fr); gap: 1em; margin: 1.5em 0; }
  .stat-card { background: #f4f4f8; border-radius: 8px; padding: 1em; text-align: center; }
  .stat-value { font-size: 1.8em; font-weight: bold; color: #6c63ff; }
  .insight-card { background: #f4f4f8; border-radius: 8px; padding: 1em; margin: 0.8em 0; }
  .insight-title { font-weight: bold; font-size: 1.1em; color: #1a1a2e; }
  .watermark { position: fixed; top: 50%; left: 50%; transform: translate(-50%,-50%) rotate(-30deg);
               font-size: 4em; color: rgba(0,0,0,0.05); font-weight: bold; pointer-events: none; z-index: 9999; }
  table { width: 100%; border-collapse: collapse; margin: 1em 0; }
  th { background: #6c63ff; color: white; padding: 0.5em; text-align: left; }
  td { padding: 0.5em; border-bottom: 1px solid #eee; }
  tr:nth-child(even) td { background: #f9f9fc; }
  .page-break { page-break-after: always; }
</style>
</head>
<body>
{% if watermarked %}<div class="watermark">ProductIQ Free Report — productiq.in</div>{% endif %}

<div class="cover page-break">
  <h1>ProductIQ</h1>
  <div style="font-size:2.2em; font-weight: bold; color: white; margin-top: 2em;">
    Product Intelligence Report
  </div>
  <div class="sub">{{ category }} | {{ brand_name or 'Market Analysis' }}</div>
  <div class="sub" style="color:#888; font-size:1em; margin-top: 1em;">{{ date or '' }}</div>
</div>

<h2>Executive Summary</h2>
{% for insight in (top_insights or insights or [])[:5] %}
<div class="insight-card">
  <div class="insight-title">{{ loop.index }}. {{ insight.title if insight is mapping else insight }}</div>
  {% if insight is mapping and insight.body %}<p>{{ insight.body[:300] }}...</p>{% endif %}
</div>
{% endfor %}

<div class="page-break"></div>
<h2>Market Overview</h2>
<div class="stat-grid">
  <div class="stat-card"><div class="stat-value">{{ products|length }}</div><div>Products Analysed</div></div>
  <div class="stat-card"><div class="stat-value">{{ competitors|length }}</div><div>Competitors Mapped</div></div>
  <div class="stat-card"><div class="stat-value">{{ review_clusters|length }}</div><div>Topic Clusters</div></div>
  <div class="stat-card"><div class="stat-value">{{ trends|length }}</div><div>Trends Tracked</div></div>
</div>

<div class="page-break"></div>
<h2>Competitive Landscape</h2>
{% if competitors %}
<table>
  <tr><th>Brand</th><th>Price (₹)</th><th>Rating</th><th>Key Strength</th></tr>
  {% for c in competitors[:10] %}
  <tr>
    <td>{{ c.brand_name }}</td>
    <td>₹{{ c.price_inr|int if c.price_inr else '—' }}</td>
    <td>{{ c.rating or '—' }}</td>
    <td>{{ (c.key_strengths or ['—'])|join(', ')|truncate(80) }}</td>
  </tr>
  {% endfor %}
</table>
{% endif %}

<div class="page-break"></div>
<h2>Product Concepts</h2>
{% for concept in (product_concepts or []) %}
<div class="insight-card">
  <div class="insight-title">{{ loop.index }}. {{ concept.concept_name }}</div>
  <p><em>{{ concept.tagline }}</em></p>
  <p><strong>Target:</strong> {{ concept.target_persona }}</p>
  <p><strong>USP:</strong> {{ concept.usp }}</p>
  <p><strong>Price:</strong> ₹{{ concept.suggested_price_inr|int if concept.suggested_price_inr else '—' }}</p>
  <p><strong>Validation Score:</strong> {{ concept.validation_score }}/100</p>
</div>
{% endfor %}

<div class="page-break"></div>
<h2>Methodology</h2>
<p>This report was generated by ProductIQ's 8-agent AI pipeline using CrewAI and Google Gemini.</p>
<p>Data sources: Amazon India, Flipkart, Google Trends, Reddit, SerpAPI.</p>
<p>Sentiment analysis: VADER + BERTopic clustering. NLP: spaCy.</p>
<p>Report generated: {{ date or 'April 2026' }}</p>
</body>
</html>"""


def _get_rfq_fallback_template() -> str:
    return """<!DOCTYPE html>
<html><head><meta charset="UTF-8"/><title>RFQ — {{ supplier.company_name }}</title>
<style>
  body { font-family: Arial, sans-serif; max-width: 800px; margin: 40px auto; color: #1a1a2e; }
  .header { border-bottom: 3px solid #6c63ff; padding-bottom: 1em; margin-bottom: 2em; }
  .logo { color: #6c63ff; font-size: 1.5em; font-weight: bold; }
  h2 { color: #6c63ff; }
  table { width: 100%; border-collapse: collapse; }
  td, th { padding: 0.5em; border: 1px solid #ddd; }
  th { background: #f4f4f8; }
</style></head>
<body>
<div class="header">
  <div class="logo">ProductIQ</div>
  <div>Date: {{ date or '—' }}</div>
</div>

<p><strong>To:</strong> {{ supplier.company_name }}<br/>
<strong>Location:</strong> {{ supplier.location or '—' }}</p>

<h2>Request for Quotation</h2>
<p><strong>Subject:</strong> RFQ for {{ concept.concept_name }}</p>

<p>Dear {{ supplier.company_name }},</p>
<p>We are developing a new product: <strong>{{ concept.concept_name }}</strong> — {{ concept.tagline }}</p>
<p>We are seeking quotations for its manufacture. Please provide:</p>

<table>
  <tr><th>Requirement</th><th>Details</th></tr>
  <tr><td>Product Category</td><td>{{ concept.gap_it_fills or '—' }}</td></tr>
  <tr><td>Required Certifications</td><td>FSSAI, ISO, GMP</td></tr>
  <tr><td>MOQ Enquiry</td><td>Please quote for 1,000 / 5,000 / 10,000 units</td></tr>
  <tr><td>Lead Time</td><td>Please indicate production + delivery timeline</td></tr>
  <tr><td>Price Per Unit</td><td>Please quote per unit at each MOQ tier</td></tr>
</table>

<p>Please respond within 5 business days. We look forward to your quotation.</p>
<p><strong>ProductIQ Sourcing Team</strong><br/>productiq.in</p>
</body></html>"""